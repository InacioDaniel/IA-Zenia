# script que permitirá ao zénia aceder e criar a arquivos no seu servidor e permitirá a ia colocar o conteudo dentro de um arquivo caso o usuario peça download do mesmo, numa fase inicial vai conseguir arquivos de codigo fonte de programação arquivos de documentos word excel powerpoint publisher txt pdf entre outros

import os
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.pdfgen import canvas
from io import BytesIO

def create_text_file(content, filename):
    with open(filename, 'w', encoding='utf-8') as f:
        f.write(content)
    return filename

def create_word_file(content, filename):
    doc = Document()
    doc.add_heading('Documento Gerado pela IA-Zénia', 0)
    doc.add_paragraph(content)
    doc.save(filename)
    return filename

def create_excel_file(data, filename):
    wb = Workbook()
    ws = wb.active
    # Se data for uma lista de listas, preenche as células
    if isinstance(data, list) and all(isinstance(row, list) for row in data):
        for row in data:
            ws.append(row)
    else:
        ws.append([data])
    wb.save(filename)
    return filename

def create_powerpoint_file(content, filename):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "IA-Zénia Presentation"
    subtitle.text = content
    prs.save(filename)
    return filename

def create_pdf_file(content, filename):
    c = canvas.Canvas(filename)
    c.drawString(100, 750, "Documento Gerado pela IA-Zénia")
    # Simples quebra de linha (limitada)
    y = 700
    for line in content.split('\n'):
        c.drawString(100, y, line)
        y -= 20
    c.save()
    return filename

def generate_file(file_type, content, filename):
    """
    Função principal para gerar arquivos baseada no tipo.
    """
    if file_type == 'txt':
        return create_text_file(content, filename)
    elif file_type == 'docx':
        return create_word_file(content, filename)
    elif file_type == 'xlsx':
        return create_excel_file(content, filename)
    elif file_type == 'pptx':
        return create_powerpoint_file(content, filename)
    elif file_type == 'pdf':
        return create_pdf_file(content, filename)
    else:
        return create_text_file(content, filename)

if __name__ == "__main__":
    # Exemplo de teste
    print("Gerando arquivos de teste...")
    generate_file('txt', 'Olá do Zénia!', 'teste.txt')
    generate_file('docx', 'Conteúdo Word do Zénia', 'teste.docx')
    print("Arquivos gerados com sucesso.")
